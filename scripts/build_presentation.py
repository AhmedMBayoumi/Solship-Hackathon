"""
Final presentation — clean, minimal, blue-and-white only.

6 slides per the official PDF:
  1. Forecasting model   (architecture diagram + RMSE/MAE/NRMSE)
  2. Controller approach (MPC architecture + H justification)
  3. March W3 dispatch plot
  4. Results table       (bills, savings, oracle gap)
  5. Generalization      (surprise NRMSE)
  6. Hardest problem & next steps
"""
from pathlib import Path
import sys, json
sys.path.insert(0, str(Path(__file__).parents[1]))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).parents[1]
OUT  = ROOT / "outputs/Solship_Hackathon_Presentation.pptx"

# Minimalist blue + white palette (only)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BLUE   = RGBColor(0x2A, 0x5D, 0xAA)   # primary blue
NAVY   = RGBColor(0x0F, 0x2D, 0x52)   # deep navy
LIGHT  = RGBColor(0xDA, 0xE5, 0xF4)   # very light blue
SOFT   = RGBColor(0xF6, 0xF8, 0xFC)   # almost-white blue tint
GREY   = RGBColor(0x55, 0x5F, 0x6B)   # body text
PALE   = RGBColor(0xE9, 0xED, 0xF2)   # divider lines

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def rect(slide, x, y, w, h, color, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if line is None: s.line.fill.background()
    else: s.line.color.rgb = line
    s.fill.solid(); s.fill.fore_color.rgb = color
    return s


def text(slide, x, y, w, h, txt, size=14, bold=False, color=NAVY,
         align=PP_ALIGN.LEFT, italic=False, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top  = Inches(0.02); tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = txt
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = font
    return tb


def title_block(slide, num, title):
    # Slim accent bar (blue), thin line under it
    rect(slide, 0, 0, Inches(0.3), SH, BLUE)
    text(slide, Inches(0.5), Inches(0.35), Inches(2.5), Inches(0.5),
         f"0{num}", size=42, bold=True, color=BLUE)
    text(slide, Inches(0.5), Inches(0.95), SW-Inches(1.0), Inches(0.6),
         title, size=24, bold=True, color=NAVY)
    rect(slide, Inches(0.5), Inches(1.55), SW-Inches(1.0), Inches(0.02), PALE)


# ─────────────────────────────────────────────────────────────────────
# SLIDE 1 — Forecasting model (visual architecture + metrics)
# ─────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, WHITE)
title_block(s, 1, "Forecasting model")

# Architecture image (full width across the top half)
s.shapes.add_picture(
    str(ROOT / "outputs/plots/presentation/architecture.png"),
    Inches(0.5), Inches(1.85), width=Inches(12.4))

# Metrics row at the bottom (4 KPI cards across)
metrics_y = Inches(5.7)
metrics_h = Inches(1.5)
metric_data = [
    ("52.17%",   "NRMSE  (★ primary)",  NAVY,  WHITE),
    ("0.470 kW", "RMSE",                 BLUE,  WHITE),
    ("0.283 kW", "MAE",                  BLUE,  WHITE),
    ("48 / 57 %","NRMSE  Apr / Sept",    BLUE,  WHITE),
]
n = len(metric_data)
gap_m = Inches(0.2)
total_w = SW - Inches(1.0)
card_w  = (total_w - (n-1)*gap_m) / n
x0 = Inches(0.5)
for i, (val, lab, fg_bg, fg_text) in enumerate(metric_data):
    x = x0 + i*(card_w + gap_m)
    rect(s, x, metrics_y, card_w, metrics_h, fg_bg)
    text(s, x, metrics_y+Inches(0.20), card_w, Inches(0.7),
         val, size=28, bold=True, color=fg_text, align=PP_ALIGN.CENTER)
    text(s, x, metrics_y+Inches(1.0), card_w, Inches(0.45),
         lab, size=11, color=fg_text, align=PP_ALIGN.CENTER)

# Small helper line under the architecture
text(s, Inches(0.5), Inches(5.40), SW-Inches(1.0), Inches(0.3),
     "2025 test set  ·  April + September  ·  5760 timesteps",
     size=10, color=GREY, italic=True, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────
# SLIDE 2 — Controller approach (MPC architecture + H justification)
# ─────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, WHITE)
title_block(s, 2, "Controller approach  ·  Rolling-horizon MPC, H = 96")

# LEFT — MPC loop diagram
mpc_x = Inches(0.7); mpc_y = Inches(1.85)
mpc_w = Inches(7.0); mpc_h = Inches(5.0)

text(s, mpc_x, mpc_y, mpc_w, Inches(0.4),
     "MPC LOOP — at each 15-min step", size=11, bold=True, color=BLUE)

step_y = mpc_y + Inches(0.5)
step_h = Inches(0.62)
step_gap = Inches(0.12)

steps = [
    ("1", "Get load forecast for next H = 96 steps (24 h)"),
    ("2", "Solve LP — 481 vars: p_chg, p_dis, p_imp, p_exp, soc"),
    ("3", "Apply hard constraints — SoC ∈ [0,1], ±8 kW battery, ±6 kW grid"),
    ("4", "Execute first action only (discard the rest)"),
    ("5", "Update SoC from realised battery action, advance to t+1"),
]
for i, (num, body) in enumerate(steps):
    y = step_y + i*(step_h + step_gap)
    rect(s, mpc_x, y, mpc_w, step_h, SOFT, line=LIGHT)
    # number circle
    rect(s, mpc_x+Inches(0.10), y+Inches(0.10), Inches(0.42), Inches(0.42), BLUE)
    text(s, mpc_x+Inches(0.10), y+Inches(0.13), Inches(0.42), Inches(0.36),
         num, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    text(s, mpc_x+Inches(0.65), y+Inches(0.16), mpc_w-Inches(0.75), Inches(0.34),
         body, size=12, color=NAVY)

# RIGHT — Solver + horizon justification
r_x = mpc_x + mpc_w + Inches(0.5)
r_w = SW - r_x - Inches(0.5)

# Solver block
text(s, r_x, Inches(1.85), r_w, Inches(0.4),
     "OPTIMIZER", size=11, bold=True, color=BLUE)
solver_y = Inches(2.25)
text(s, r_x, solver_y, r_w, Inches(0.4),
     "scipy.optimize.linprog  (HiGHS)",
     size=15, bold=True, color=NAVY)
text(s, r_x, solver_y+Inches(0.45), r_w, Inches(1.6),
     "→  Linear program — convex, polynomial-time, single global optimum\n"
     "→  Open-source, deterministic, ~10 ms per H=96 solve\n"
     "→  No MIP needed: at non-negative tariffs, LP picks one direction naturally",
     size=11, color=GREY)

# Horizon table
h_y = Inches(4.5)
text(s, r_x, h_y, r_w, Inches(0.4),
     "WHY H = 96  (24-hour horizon)",
     size=11, bold=True, color=BLUE)

table_top = h_y + Inches(0.5)
rows = [("H", "Bill (€)", "vs A"),
        ("4",  "+63.71", "+71"),
        ("24", "+6.22",  "+14"),
        ("48", "−16.71", "−9"),
        ("96", "−18.89", "−11"),
       ]
col_xs = [r_x+Inches(0.10), r_x+r_w*0.32, r_x+r_w*0.62]
col_ws = [r_w*0.30, r_w*0.30, r_w*0.30]
for i, row in enumerate(rows):
    yy = table_top + Inches(i*0.30)
    is_h = (i == 0); is_w = (i == len(rows)-1)
    if is_h: rect(s, r_x, yy, r_w, Inches(0.28), NAVY)
    if is_w: rect(s, r_x, yy, r_w, Inches(0.28), LIGHT)
    fg = WHITE if is_h else NAVY
    for c, cell in enumerate(row):
        text(s, col_xs[c], yy+Inches(0.02), col_ws[c], Inches(0.26),
             cell, size=11, bold=is_h or is_w, color=fg, align=PP_ALIGN.CENTER)

text(s, r_x, table_top+Inches(len(rows)*0.30+0.18), r_w, Inches(0.5),
     "Captures full F3-night → F1-day arbitrage cycle.",
     size=10, color=GREY, italic=True)


# ─────────────────────────────────────────────────────────────────────
# SLIDE 3 — March Week 3 dispatch plot (mandatory)
# ─────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, WHITE)
title_block(s, 3, "March Week 3 dispatch  ·  Mar 17 – 23, 2025")

s.shapes.add_picture(
    str(ROOT / "outputs/plots/presentation/march_week3_dispatch.png"),
    Inches(0.7), Inches(1.85), width=Inches(11.9))

text(s, Inches(0.7), Inches(7.05), Inches(11.9), Inches(0.35),
     "All hard constraints respected: battery [-5.2, +3.8] kW (limit ±8) · grid [-1.4, +4.7] kW (limit ±6) · SoC [0%, 97.5%]",
     size=10, color=GREY, italic=True, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────
# SLIDE 4 — Results
# ─────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, WHITE)
title_block(s, 4, "Results  ·  90% of oracle savings  ·  +154% vs Baseline A")

# 4-line plot left
s.shapes.add_picture(
    str(ROOT / "outputs/plots/presentation/cumulative_bill_4line.png"),
    Inches(0.5), Inches(1.85), width=Inches(8.0))
text(s, Inches(0.5), Inches(5.0), Inches(8.0), Inches(0.3),
     "Cumulative electricity bill — April + September 2025",
     size=10, color=GREY, italic=True, align=PP_ALIGN.CENTER)

# Bills table on right
bills = json.loads((ROOT / "outputs/models/presentation_bills.json").read_text())
r_x = Inches(8.7); r_w = SW - r_x - Inches(0.5)
text(s, r_x, Inches(1.85), r_w, Inches(0.4),
     "Bill summary  (Apr + Sept, H=96)", size=12, bold=True, color=BLUE)

rows = [
    ("Method",                    "Total"),
    ("Baseline B (no battery)",   f"€{bills['baseline_B']['total']:+.2f}"),
    ("Baseline A (historical)",   f"€{bills['baseline_A']['total']:+.2f}"),
    ("Our controller (MPC)",      f"€{bills['ours_H96']['total']:+.2f}"),
    ("Oracle (perfect)",          f"€{bills['oracle_H96']['total']:+.2f}"),
]
table_y = Inches(2.4)
for i, (k, v) in enumerate(rows):
    yy = table_y + Inches(i*0.42)
    is_h = (i == 0); is_us = (i == 3)
    if is_h:
        rect(s, r_x, yy, r_w, Inches(0.36), NAVY)
        fg = WHITE
    elif is_us:
        rect(s, r_x, yy, r_w, Inches(0.36), LIGHT)
        fg = NAVY
    else:
        fg = NAVY
    text(s, r_x+Inches(0.15), yy+Inches(0.05), r_w*0.62, Inches(0.30),
         k, size=11, bold=is_h or is_us, color=fg)
    text(s, r_x+r_w*0.62, yy+Inches(0.05), r_w*0.32, Inches(0.30),
         v, size=12, bold=is_h or is_us, color=fg, align=PP_ALIGN.RIGHT)

# Savings + oracle gap
sav = bills["baseline_A"]["total"] - bills["ours_H96"]["total"]
sav_pct = sav / abs(bills["baseline_A"]["total"]) * 100
gap = bills["ours_H96"]["total"] - bills["oracle_H96"]["total"]
captured = (bills["baseline_A"]["total"] - bills["ours_H96"]["total"]) / \
           (bills["baseline_A"]["total"] - bills["oracle_H96"]["total"]) * 100

kpi_y = Inches(5.0)
text(s, r_x, kpi_y, r_w, Inches(0.4),
     "Savings vs Baseline A", size=11, color=GREY)
text(s, r_x, kpi_y+Inches(0.35), r_w, Inches(0.6),
     f"+€{sav:.2f}  ({sav_pct:+.0f}%)", size=22, bold=True, color=NAVY)

text(s, r_x, kpi_y+Inches(1.05), r_w, Inches(0.4),
     "Oracle gap", size=11, color=GREY)
text(s, r_x, kpi_y+Inches(1.40), r_w, Inches(0.6),
     f"€{gap:.2f}   ({captured:.0f}% captured)", size=22, bold=True, color=NAVY)

text(s, r_x, Inches(7.05), r_w, Inches(0.3),
     "✓ Extension 1 done (H sweep 4 → 96)", size=10, color=BLUE, italic=True)


# ─────────────────────────────────────────────────────────────────────
# SLIDE 5 — Generalization (surprise dataset)
# ─────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, WHITE)
title_block(s, 5, "Generalization  ·  surprise dataset (March 2026)")

s.shapes.add_picture(
    str(ROOT / "outputs/plots/surprise/surprise_forecast_full.png"),
    Inches(0.7), Inches(1.85), width=Inches(8.5))
text(s, Inches(0.7), Inches(4.20), Inches(8.5), Inches(0.3),
     "March 2026 surprise dataset — full month",
     size=10, color=GREY, italic=True, align=PP_ALIGN.CENTER)

# Big NRMSE on right
r_x = Inches(9.4); r_w = SW - r_x - Inches(0.5)
text(s, r_x, Inches(1.85), r_w, Inches(0.4),
     "Surprise dataset NRMSE", size=11, bold=True, color=BLUE)
text(s, r_x, Inches(2.4), r_w, Inches(1.2),
     "32.16%", size=72, bold=True, color=NAVY)
text(s, r_x, Inches(3.65), r_w, Inches(0.4),
     "down from 52.17% on 2025", size=11, color=GREY, italic=True)

text(s, r_x, Inches(4.4), r_w, Inches(0.4),
     "RMSE",  size=11, color=GREY)
text(s, r_x, Inches(4.7), r_w, Inches(0.4),
     "1.05 kW", size=14, bold=True, color=NAVY)
text(s, r_x, Inches(5.2), r_w, Inches(0.4),
     "MAE",  size=11, color=GREY)
text(s, r_x, Inches(5.5), r_w, Inches(0.4),
     "0.63 kW", size=14, bold=True, color=NAVY)

text(s, Inches(0.7), Inches(6.6), SW-Inches(1.4), Inches(0.5),
     "Same architecture — no retuning, just rebuilt features and refit on the new site (3× larger consumer → better signal-to-noise).",
     size=12, color=GREY, italic=True, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────
# SLIDE 6 — Hardest problem & next steps
# ─────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, WHITE)
title_block(s, 6, "Hardest problem  ·  Next steps")

# Two columns
col_w = (SW - Inches(1.4) - Inches(0.4))/2
left_x = Inches(0.7); right_x = left_x + col_w + Inches(0.4)
top = Inches(1.85)

# LEFT — hardest problem
text(s, left_x, top, col_w, Inches(0.4),
     "HARDEST PROBLEM", size=11, bold=True, color=BLUE)
text(s, left_x, top+Inches(0.45), col_w, Inches(0.5),
     "NRMSE ≠ Bill", size=22, bold=True, color=NAVY)

text(s, left_x, top+Inches(1.1), col_w, Inches(2.5),
     "We discovered a Pareto trade-off:\n"
     "smoothing the forecast lowered NRMSE\n"
     "(60.4 % → 52.2 %), but the LP under-charged\n"
     "before peaks, raising the bill\n"
     "(−€18.89 → −€18.23).\n\n"
     "NRMSE penalises errors symmetrically,\n"
     "but the bill is asymmetric in peak timing.",
     size=12, color=GREY)

text(s, left_x, top+Inches(3.6), col_w, Inches(0.4),
     "FIX  —  spike-preserving coring", size=11, bold=True, color=BLUE)
text(s, left_x, top+Inches(4.0), col_w, Inches(1.5),
     "Smooth where |raw − MA(3)| < 0.4 kW (= noise),\n"
     "keep raw where |residual| > 0.4 kW (= real spike).\n\n"
     "→ Pareto win: NRMSE 60.4 → 57.8 %\n"
     "    AND bill −€18.89 → −€18.91",
     size=11, color=NAVY)

# RIGHT — next steps
text(s, right_x, top, col_w, Inches(0.4),
     "NEXT STEPS", size=11, bold=True, color=BLUE)
text(s, right_x, top+Inches(0.45), col_w, Inches(0.5),
     "Given one more day", size=22, bold=True, color=NAVY)

next_steps = [
    ("Decision-focused loss",
     "Train the forecaster directly on the bill (SPO+ — Elmachtoub & Grigas, Mgmt Sci 2022)."),
    ("PV forecast fusion",
     "Pair load forecast with a separate PV model using ARPA radiation directly."),
    ("Tariff-band-aware smoothing",
     "Smooth only F3 hours; keep F1 raw for peak fidelity."),
    ("Battery degradation cost",
     "Add cycle-count penalty to LP (already wired, set to 0 here)."),
]
for i, (h, b) in enumerate(next_steps):
    yy = top + Inches(1.1 + i*1.05)
    rect(s, right_x, yy, Inches(0.06), Inches(0.85), BLUE)
    text(s, right_x+Inches(0.20), yy, col_w-Inches(0.20), Inches(0.32),
         h, size=12, bold=True, color=NAVY)
    text(s, right_x+Inches(0.20), yy+Inches(0.32), col_w-Inches(0.20), Inches(0.55),
         b, size=10.5, color=GREY)

# Footer
rect(s, 0, Inches(7.05), SW, Inches(0.45), NAVY)
text(s, Inches(0.5), Inches(7.13), SW-Inches(1.0), Inches(0.3),
     "Ahmed Mohamed Bayoumi  ·  Team 25  ·  Energy AI Hackathon 2026",
     size=11, color=WHITE, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"Saved -> {OUT}")
print(f"Slides: {len(prs.slides)}")
print(f"Size:   {OUT.stat().st_size/1024:.0f} KB")
